from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
RESULTS_DIR = BASE_DIR / "results"
OUTPUT_PATH = BASE_DIR / "Brain_Tumor_Project_Presentation.pptx"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()

    for idx, item in enumerate(bullets):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
        else:
            p.text = str(item)
            p.level = 0
        p.font.size = Pt(22 if p.level == 0 else 18)


def add_image_slide(prs, title, image_path, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    if image_path.exists():
        left = Inches(0.6)
        top = Inches(1.3)
        width = Inches(12.1)
        height = Inches(5.5)
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.5))
        tf = tx_box.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(14)
    else:
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(2.5))
        tf = tx_box.text_frame
        tf.text = f"Image not found: {image_path.name}"
        tf.paragraphs[0].font.size = Pt(24)


def add_comparison_table_slide(prs, title, csv_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    if not csv_path.exists():
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(2.5))
        tf = tx_box.text_frame
        tf.text = f"Comparison file not found: {csv_path.name}"
        tf.paragraphs[0].font.size = Pt(24)
        return

    df = pd.read_csv(csv_path)
    cols = ["model", "accuracy", "checkpoint"]
    for col in cols:
        if col not in df.columns:
            tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(2.5))
            tf = tx_box.text_frame
            tf.text = "model_comparison.csv format is unexpected."
            tf.paragraphs[0].font.size = Pt(24)
            return

    display_df = df[cols].copy()
    display_df["accuracy"] = display_df["accuracy"].map(lambda x: f"{x:.4f}")

    rows = len(display_df) + 1
    cols_n = len(display_df.columns)
    table = slide.shapes.add_table(rows, cols_n, Inches(1.0), Inches(1.8), Inches(11.0), Inches(2.4)).table

    for col_idx, col_name in enumerate(display_df.columns):
        table.cell(0, col_idx).text = col_name.upper()
        for run in table.cell(0, col_idx).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)

    for row_idx, row in enumerate(display_df.itertuples(index=False), start=1):
        for col_idx, val in enumerate(row):
            table.cell(row_idx, col_idx).text = str(val)
            table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(15)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Brain Tumor Classification from MRI",
        "Vision Transformer vs ResNet50 | End-to-End Research + Deployment Work",
    )

    add_bullets_slide(
        prs,
        "Problem Statement",
        [
            "Early and accurate brain tumor classification from MRI is clinically important.",
            "Goal: 4-class classification (glioma, meningioma, notumor, pituitary).",
            "Need: research-grade comparison + practical deployment-ready system.",
        ],
    )

    add_bullets_slide(
        prs,
        "Project Objectives",
        [
            "Build reproducible deep learning pipeline in PyTorch.",
            "Benchmark ResNet50 and ViT-B/16 under same train/eval setup.",
            "Deploy model via FastAPI and Streamlit GUI.",
            "Enable continual learning from newly labeled uploaded scans.",
        ],
    )

    add_bullets_slide(
        prs,
        "Dataset Summary",
        [
            "Total MRI images: 7,023",
            "Training: 5,712 | Testing: 1,311",
            "Classes: glioma (1621), meningioma (1645), notumor (2000), pituitary (1757)",
            "Image size: 224x224, train-time augmentation + model-specific preprocessing",
        ],
    )

    add_bullets_slide(
        prs,
        "Methodology Pipeline",
        [
            "Data loading: ImageFolder + DataLoader from src/data_loader.py",
            "Model options: ResNet50 / ViT-B-16 from src/model.py",
            "Training: CrossEntropyLoss + Adam (lr=1e-4), checkpoint saving",
            "Evaluation: classification report + confusion matrix + accuracy",
            "Experiment orchestration: run_experiments.py",
        ],
    )

    add_bullets_slide(
        prs,
        "Model Architectures",
        [
            "ResNet50: ImageNet pretrained backbone, FC replaced for 4 classes",
            "ViT-B/16: ImageNet pretrained transformer, head replaced for 4 classes",
            "Unified interface allows fair side-by-side comparison",
            "Selected model and checkpoint metadata saved for inference",
        ],
    )

    add_image_slide(
        prs,
        "Learning Curve (ResNet)",
        RESULTS_DIR / "accuracy_loss_resnet.png",
        "Train/validation accuracy-loss trajectory for ResNet",
    )

    add_image_slide(
        prs,
        "Learning Curve (ViT)",
        RESULTS_DIR / "accuracy_loss_vit.png",
        "Train/validation accuracy-loss trajectory for ViT",
    )

    add_image_slide(
        prs,
        "Confusion Matrix (ResNet)",
        RESULTS_DIR / "confusion_matrix_resnet.png",
        "Class-wise prediction behavior for ResNet model",
    )

    add_image_slide(
        prs,
        "Confusion Matrix (ViT)",
        RESULTS_DIR / "confusion_matrix_vit.png",
        "Class-wise prediction behavior for ViT model",
    )

    add_comparison_table_slide(
        prs,
        "Model Comparison",
        RESULTS_DIR / "model_comparison.csv",
    )

    add_bullets_slide(
        prs,
        "Key Results",
        [
            "ResNet accuracy: ~0.9703 (best in current comparison)",
            "ViT accuracy: ~0.9207",
            "Both models are deployment-ready through shared inference stack",
            "Artifacts generated: reports, confusion matrices, learning curves, comparison table",
        ],
    )

    add_bullets_slide(
        prs,
        "Deployment and Continual Learning",
        [
            "FastAPI endpoint: /predict for programmatic inference",
            "Streamlit app: image upload, prediction, admin review and labeling",
            "Labeled uploads can trigger fine-tuning via continual_train.py",
            "Inference supports local checkpoint fallback and cloud checkpoint download",
        ],
    )

    add_bullets_slide(
        prs,
        "Project Structure (Highlights)",
        [
            "src/: config, model, data_loader, train, evaluate, continual_train, gradcam",
            "deployment/: inference.py, app.py, gui_app.py",
            "results/: model metrics and visual artifacts",
            "models/: best_model_resnet.pth, best_model_vit.pth",
            "main.py and run_experiments.py as primary entry points",
        ],
    )

    add_bullets_slide(
        prs,
        "Limitations and Future Work",
        [
            "No cross-validation pipeline yet",
            "Statistical confidence intervals across repeated runs are pending",
            "Grad-CAM currently focused on ResNet",
            "Future: richer explainability, domain generalization, clinical metadata integration",
        ],
    )

    add_bullets_slide(
        prs,
        "Conclusion",
        [
            "Successfully built an end-to-end MRI tumor classification system.",
            "Compared transformer and CNN backbones under one reproducible setup.",
            "Delivered training-to-deployment workflow with continual improvement loop.",
            "Current best performer in this run: ResNet50.",
        ],
    )

    add_bullets_slide(
        prs,
        "Thank You",
        [
            "Brain Tumor Classification using Deep Learning",
            "Prepared from project artifacts and implementation outputs",
            "Author: Abhishek Gautam",
        ],
    )

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    out = build_presentation()
    print(f"Presentation created: {out}")
